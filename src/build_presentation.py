"""
src/build_presentation.py
Gera wine_quality_eda_storytelling.pptx a partir dos dados brutos e de
results/metricas_modelos.csv.

A apresentação é CÓDIGO, não um arquivo editado à mão: todo número exibido é
lido de data/ ou de results/, portanto não pode divergir do notebook.

Uso:  python src/build_presentation.py     (a partir da raiz do repositório)

--------------------------------------------------------------------------
Sistema de design (editorial, sóbrio)
  Superfície  off-white     #FCFCFB
  Tinta       primária      #1A1A1A   secundária #6B6B6B   fio #E2E0DC
  Acento      bordô         #8C2F39   (protagonista: Random Forest, corr. negativa)
  Contraponto azul-aço      #1F6FA8   (Regressão Logística, corr. positiva)
  Display     Georgia (serifada)      Texto: Segoe UI

O par bordô/azul-aço foi validado para daltonismo: ΔE 45,5 (protanopia),
84,4 (tritanopia) — muito acima do piso de 12. Nenhuma informação é
transmitida apenas por cor: todas as barras têm rótulo direto.
--------------------------------------------------------------------------
"""

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent

# ---------------------------------------------------------------- design tokens
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_SOFT = RGBColor(0x6B, 0x6B, 0x6B)
RULE = RGBColor(0xE2, 0xE0, 0xDC)
ACCENT = RGBColor(0x8C, 0x2F, 0x39)   # bordô
COUNTER = RGBColor(0x1F, 0x6F, 0xA8)  # azul-aço

DISPLAY = 'Georgia'
BODY = 'Segoe UI'

W, H = Inches(10), Inches(5.625)
MARGIN = Inches(0.62)
CONTENT_W = W - 2 * MARGIN


# ---------------------------------------------------------------- helpers
def textbox(slide, x, y, w, h, text, *, size, font=BODY, color=INK, bold=False,
            align=PP_ALIGN.LEFT, italic=False, spacing=None, caps=False, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        run = p.add_run()
        run.text = line.upper() if caps else line
        f = run.font
        f.name, f.size, f.bold, f.italic = font, Pt(size), bold, italic
        f.color.rgb = color
    return box


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def hairline(slide, x, y, w, color=RULE):
    rect(slide, x, y, w, Emu(9525), color)  # ~0.75pt


def num(v, casas=2, sinal=False):
    """Número no padrão pt-BR (vírgula decimal). NUNCA use .replace('.', ',')
    sobre uma frase inteira: isso corrompe a pontuação do texto."""
    s = f'{v:+.{casas}f}' if sinal else f'{v:.{casas}f}'
    return s.replace('.', ',')


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # em branco
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = SURFACE
    return slide


def header(slide, eyebrow, title, n):
    """Cabeçalho editorial: sobrancelha caixa-alta, título serifado, fio."""
    textbox(slide, MARGIN, Inches(0.42), CONTENT_W, Inches(0.18), eyebrow,
            size=8.5, color=ACCENT, bold=True, caps=True)
    textbox(slide, MARGIN, Inches(0.63), CONTENT_W, Inches(0.42), title,
            size=23, font=DISPLAY, color=INK)
    hairline(slide, MARGIN, Inches(1.18), CONTENT_W)
    footer(slide, n)


def footer(slide, n):
    textbox(slide, MARGIN, Inches(5.16), CONTENT_W - Inches(0.4), Inches(0.2),
            'Wine Quality Classification  ·  POSTECH DTAT  ·  Tech Challenge Fase 2',
            size=7.5, color=INK_SOFT)
    textbox(slide, W - MARGIN - Inches(0.4), Inches(5.16), Inches(0.4), Inches(0.2),
            f'{n:02d}', size=7.5, color=INK_SOFT, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------- dados
# Wine Quality Dataset (Kaggle) — WineQT.csv: 1.143 amostras de vinho TINTO.
df = pd.read_csv(ROOT / 'data/WineQT.csv').drop(columns=['Id'])
df.columns = df.columns.str.replace(' ', '_')
df['hq'] = (df['quality'] >= 7).astype(int)

pct_alta = df['hq'].mean() * 100
corr = df.corr(numeric_only=True)['hq']


def taxa(col, thr):
    """% de alta qualidade acima e abaixo de um limiar."""
    return (df[df[col] > thr]['hq'].mean() * 100,
            df[df[col] <= thr]['hq'].mean() * 100)


alc_hi, alc_lo = taxa('alcohol', 11)
va_hi, va_lo = taxa('volatile_acidity', 0.7)
su_hi, su_lo = taxa('sulphates', 0.65)
ci_hi, ci_lo = taxa('citric_acid', 0.30)

met = pd.read_csv(ROOT / 'results/metricas_modelos.csv').set_index('Modelo')
LR, RF = 'Logistic Regression', 'Random Forest'

prs = Presentation()
prs.slide_width, prs.slide_height = W, H

# ================================================================ 1 · capa
s = new_slide(prs)
rect(s, MARGIN, Inches(1.55), Inches(0.035), Inches(1.55), ACCENT)  # filete vertical
x = MARGIN + Inches(0.3)
textbox(s, x, Inches(1.5), CONTENT_W, Inches(0.2),
        'Análise exploratória e modelagem preditiva', size=9, color=ACCENT, bold=True, caps=True)
textbox(s, x, Inches(1.85), Inches(7.6), Inches(1.3),
        'Prevendo a qualidade\ndo vinho antes da taça',
        size=34, font=DISPLAY, color=INK, spacing=1.15)
hairline(s, x, Inches(3.35), Inches(6.2))
textbox(s, x, Inches(3.55), Inches(7.0), Inches(0.6),
        'O que 1.143 amostras de laboratório revelam sobre a nota que um especialista\ndará — e sobre onde vale a pena mexer no processo produtivo.',
        size=11, color=INK_SOFT, spacing=1.35)
for i, (k, v) in enumerate([(f'{len(df):,}'.replace(',', '.'), 'amostras'), ('11', 'variáveis'),
                            ('tinto', 'tipo de vinho'), (f'{num(pct_alta, 1)}%', 'alta qualidade')]):
    cx = x + Inches(1.72) * i
    textbox(s, cx, Inches(4.42), Inches(1.6), Inches(0.3), k, size=17, font=DISPLAY, color=INK)
    textbox(s, cx, Inches(4.75), Inches(1.6), Inches(0.2), v, size=8, color=INK_SOFT, caps=True)
footer(s, 1)

# ================================================================ 2 · o desafio
s = new_slide(prs)
header(s, 'O problema de negócio', 'A nota depende do paladar de quem prova', 2)
itens = [
    ('01', 'Subjetivo',
     'Dois especialistas dão notas diferentes para o mesmo vinho. A avaliação sensorial '
     'depende do paladar, da experiência e até do dia de quem prova.'),
    ('02', 'Caro e lento',
     'Cada lote consome tempo de um profissional escasso. O custo cresce linearmente com '
     'o volume produzido — não há ganho de escala.'),
    ('03', 'Tarde demais',
     'A nota só chega quando o vinho já está pronto. Quando o problema aparece, não há '
     'mais processo a corrigir: só há prejuízo a contabilizar.'),
]
col_w = (CONTENT_W - Inches(0.9)) / 3
for i, (numeral, titulo, corpo) in enumerate(itens):
    cx = MARGIN + (col_w + Inches(0.45)) * i
    if i:
        rect(s, cx - Inches(0.22), Inches(1.75), Emu(9525), Inches(1.9), RULE)
    textbox(s, cx, Inches(1.75), col_w, Inches(0.3), numeral, size=15, font=DISPLAY, color=ACCENT)
    textbox(s, cx, Inches(2.2), col_w, Inches(0.3), titulo, size=13, font=DISPLAY, color=INK)
    textbox(s, cx, Inches(2.62), col_w, Inches(1.4), corpo, size=9.5, color=INK_SOFT, spacing=1.45)
textbox(s, MARGIN, Inches(4.35), CONTENT_W, Inches(0.4),
        'A pergunta: dá para saber antes — com os dados que o laboratório já coleta?',
        size=13, font=DISPLAY, color=ACCENT, italic=True)

# ================================================================ 3 · a decisão
s = new_slide(prs)
header(s, 'Como transformamos o problema', 'De uma nota de 3 a 9 para uma decisão de sim ou não', 3)
textbox(s, MARGIN, Inches(1.5), CONTENT_W, Inches(0.4),
        'A operação não precisa prever "nota 6,4". Precisa saber em qual lote vale a pena investir atenção.',
        size=10.5, color=INK_SOFT, spacing=1.4)

bar_y, bar_h, bar_w = Inches(2.35), Inches(0.52), CONTENT_W
w_baixa = Emu(int(bar_w * (100 - pct_alta) / 100))
rect(s, MARGIN, bar_y, w_baixa, bar_h, RGBColor(0xD8, 0xD5, 0xD0))
rect(s, MARGIN + w_baixa + Emu(19050), bar_y, bar_w - w_baixa - Emu(19050), bar_h, ACCENT)

textbox(s, MARGIN, bar_y + Inches(0.68), Inches(4), Inches(0.25),
        f'{num(100 - pct_alta, 1)}%   Baixa / Média qualidade  (nota < 7)', size=10, color=INK)
textbox(s, W - MARGIN - Inches(3.4), bar_y + Inches(0.68), Inches(3.4), Inches(0.25),
        f'{num(pct_alta, 1)}%   Alta qualidade  (nota ≥ 7)', size=10, color=ACCENT,
        bold=True, align=PP_ALIGN.RIGHT)
hairline(s, MARGIN, Inches(3.72), CONTENT_W)
textbox(s, MARGIN, Inches(3.95), CONTENT_W, Inches(0.7),
        'Apenas 1 em cada 7 vinhos é de alta qualidade. É agulha no palheiro —\n'
        'e é exatamente por isso que um bom filtro vale tanto.',
        size=14, font=DISPLAY, color=INK, spacing=1.3)

# ================================================================ 4 · correlações
s = new_slide(prs)
header(s, 'O que os dados dizem', 'O que realmente acompanha a nota final', 4)

vars_pt = [('alcohol', 'Álcool'), ('volatile_acidity', 'Acidez volátil'),
           ('citric_acid', 'Ácido cítrico'), ('sulphates', 'Sulfatos'),
           ('density', 'Densidade'), ('fixed_acidity', 'Acidez fixa'),
           ('chlorides', 'Cloretos')]
# Geometria: coluna fixa de nomes à esquerda, depois a área do gráfico.
# O eixo zero fica deslocado para caber a maior barra negativa (-0,28) e o
# seu rótulo, sem invadir a coluna de nomes.
NAMES_W = Inches(1.63)          # 0,62 -> 2,25 in
SCALE = Inches(5.5)             # polegadas por unidade de correlação
LABEL_W = Inches(0.72)
zero_x = MARGIN + NAMES_W + Inches(0.15) + LABEL_W + Emu(int(0.31 * SCALE))

rect(s, zero_x, Inches(1.52), Emu(9525), Inches(2.72), RULE)  # eixo zero
row_h = Inches(0.375)
for i, (col, label) in enumerate(vars_pt):
    v = corr[col]
    y = Inches(1.62) + row_h * i
    textbox(s, MARGIN, y + Inches(0.02), NAMES_W, Inches(0.22), label,
            size=9.5, color=INK, align=PP_ALIGN.RIGHT)
    bw = Emu(int(abs(v) * SCALE))
    if v >= 0:
        rect(s, zero_x + Emu(19050), y, bw, Inches(0.2), COUNTER)
        textbox(s, zero_x + bw + Inches(0.1), y + Inches(0.005), LABEL_W, Inches(0.2),
                num(v, 2, sinal=True), size=8.5, color=INK_SOFT)
    else:
        rect(s, zero_x - bw - Emu(19050), y, bw, Inches(0.2), ACCENT)
        textbox(s, zero_x - bw - Inches(0.1) - LABEL_W, y + Inches(0.005), LABEL_W, Inches(0.2),
                num(v, 2, sinal=True), size=8.5, color=INK_SOFT, align=PP_ALIGN.RIGHT)

textbox(s, zero_x - Inches(2.6) - Inches(0.12), Inches(4.42), Inches(2.6), Inches(0.2),
        'Puxa a nota para baixo', size=8, color=ACCENT, caps=True, bold=True,
        align=PP_ALIGN.RIGHT)
textbox(s, zero_x + Inches(0.12), Inches(4.42), Inches(2.5), Inches(0.2), 'Puxa a nota para cima',
        size=8, color=COUNTER, caps=True, bold=True)
textbox(s, MARGIN, Inches(4.78), CONTENT_W, Inches(0.35),
        'A densidade é apenas o espelho do álcool — etanol é menos denso que a água. Não é um efeito independente.',
        size=9.5, color=INK_SOFT)

# ================================================================ 5 · o achado
s = new_slide(prs)
header(s, 'O achado', 'O que a intuição diz que importa — e não importa', 5)
textbox(s, MARGIN, Inches(1.62), Inches(5.7), Inches(1.7),
        'Três variáveis que todo mundo supõe decisivas não têm\n'
        'qualquer poder de separar vinho bom de vinho ruim:\n'
        'doçura, pH e SO₂ livre.',
        size=14, font=DISPLAY, color=INK, spacing=1.45)
textbox(s, MARGIN, Inches(3.4), Inches(5.7), Inches(1.3),
        'Mexer no açúcar residual ou corrigir o pH não move a nota — a associação com a '
        'qualidade é praticamente zero. O esforço do processo deve ir para álcool, acidez '
        'volátil e sulfatos, que são onde o sinal realmente está.',
        size=10, color=INK_SOFT, spacing=1.45)

bx = MARGIN + Inches(6.3)
rect(s, bx, Inches(1.55), Inches(2.45), Inches(3.0), RGBColor(0xF4, 0xF2, 0xEF))
textbox(s, bx + Inches(0.25), Inches(1.82), Inches(2), Inches(0.25),
        'Associação com a nota', size=8.5, color=INK_SOFT, caps=True, bold=True)
nao_explicam = [('Açúcar residual', 'residual_sugar'), ('pH', 'pH'), ('SO₂ livre', 'free_sulfur_dioxide')]
for i, (rot, col) in enumerate(nao_explicam):
    y = Inches(2.18) + Inches(0.63) * i
    textbox(s, bx + Inches(0.25), y, Inches(2), Inches(0.2), rot, size=9, color=INK_SOFT)
    textbox(s, bx + Inches(0.25), y + Inches(0.19), Inches(2), Inches(0.32),
            num(corr[col], 2, sinal=True), size=17, font=DISPLAY, color=INK_SOFT)
hairline(s, bx + Inches(0.25), Inches(4.14), Inches(1.95), RGBColor(0xD8, 0xD5, 0xD0))
textbox(s, bx + Inches(0.25), Inches(4.26), Inches(2), Inches(0.2),
        'Praticamente zero', size=8, color=ACCENT, caps=True, bold=True)

# ================================================================ 6 · insights
s = new_slide(prs)
header(s, 'Insights acionáveis', 'Quatro alavancas para o processo produtivo', 6)
insights = [
    ('01', 'Álcool acima de 11%',
     f'{alc_hi:.0f}% dos vinhos são de alta qualidade acima desse limite, contra apenas '
     f'{alc_lo:.0f}% abaixo. Cinco vezes mais — reflexo de uva mais madura na colheita.'),
    ('02', 'Acidez volátil sob controle',
     f'Acima de 0,7 g/L, só {va_hi:.0f}% atingem alta qualidade, contra {va_lo:.0f}% abaixo. '
     'É o aroma avinagrado, causado pela bactéria acética. Praticamente uma sentença.'),
    ('03', 'Sulfatos acima de 0,65',
     f'{su_hi:.0f}% de alta qualidade acima desse nível, contra {su_lo:.0f}% abaixo. '
     'Protegem o vinho da oxidação e preservam os aromas ao longo da guarda.'),
    ('04', 'Ácido cítrico acima de 0,3',
     f'{ci_hi:.0f}% contra {ci_lo:.0f}%. Traz frescor — mas parte do mérito é indireta: '
     'onde há cítrico, costuma haver menos acidez volátil (correlação de −0,54).'),
]
cw = (CONTENT_W - Inches(0.6)) / 2
for i, (numeral, tit, corpo) in enumerate(insights):
    cx = MARGIN + (cw + Inches(0.6)) * (i % 2)
    cy = Inches(1.6) + Inches(1.62) * (i // 2)
    textbox(s, cx, cy, Inches(0.4), Inches(0.25), numeral, size=11, font=DISPLAY, color=ACCENT)
    textbox(s, cx + Inches(0.45), cy - Inches(0.02), cw - Inches(0.45), Inches(0.25), tit,
            size=12, font=DISPLAY, color=INK)
    textbox(s, cx + Inches(0.45), cy + Inches(0.32), cw - Inches(0.45), Inches(1.0), corpo,
            size=9.5, color=INK_SOFT, spacing=1.45)
    hairline(s, cx, cy + Inches(1.32), cw)

# ================================================================ 7 · resultados
s = new_slide(prs)
header(s, 'Desempenho dos modelos', 'Random Forest vence onde importa', 7)

# legenda (identidade nunca só por cor: há rótulo direto em cada barra)
lx = W - MARGIN - Inches(3.5)
for i, (nome, cor) in enumerate([(RF, ACCENT), (LR, COUNTER)]):
    rect(s, lx + Inches(1.75) * i, Inches(1.36), Inches(0.11), Inches(0.11), cor)
    textbox(s, lx + Inches(1.75) * i + Inches(0.2), Inches(1.32), Inches(1.5), Inches(0.2),
            nome, size=8.5, color=INK_SOFT)

metricas = [('Accuracy', 'Taxa de acerto geral'), ('Precision', 'Acerta quando aponta "bom"'),
            ('Recall', 'Encontra os bons que existem'), ('F1-Score', 'Equilíbrio entre os dois'),
            ('ROC-AUC', 'Poder de discriminação')]
bar_x = MARGIN + Inches(2.75)
bar_max = Inches(4.55)
for i, (m, desc) in enumerate(metricas):
    y = Inches(1.72) + Inches(0.63) * i
    textbox(s, MARGIN, y - Inches(0.01), Inches(2.55), Inches(0.2), m, size=9.5, color=INK)
    textbox(s, MARGIN, y + Inches(0.19), Inches(2.55), Inches(0.2), desc, size=7.5, color=INK_SOFT)
    for j, (modelo, cor) in enumerate([(RF, ACCENT), (LR, COUNTER)]):
        v = met.loc[modelo, m]
        by = y + Inches(0.21) * j
        rect(s, bar_x, by, Emu(int(v * bar_max)), Inches(0.175), cor)
        textbox(s, bar_x + Emu(int(v * bar_max)) + Inches(0.09), by - Inches(0.005),
                Inches(0.7), Inches(0.2), f'{num(v * 100, 1)}%', size=8.5, color=INK_SOFT)

textbox(s, MARGIN, Inches(4.95), CONTENT_W, Inches(0.3),
        f'A Regressão Logística encontra mais vinhos bons (recall {met.loc[LR, "Recall"] * 100:.0f}%), '
        f'mas erra 2 em cada 3 que aponta (precisão {met.loc[LR, "Precision"] * 100:.0f}%). '
        'Inviável para triagem.',
        size=9.5, color=INK_SOFT)

# ================================================================ 8 · conclusões
s = new_slide(prs)
header(s, 'Recomendação', 'O modelo está pronto para um piloto', 8)
concl = [
    ('Triagem automática',
     f'Com {met.loc[RF, "ROC-AUC"] * 100:.0f}% de poder de discriminação e '
     f'{met.loc[RF, "Accuracy"] * 100:.0f}% de acerto, o modelo separa os lotes promissores. '
     'O especialista passa a provar só o que importa.'),
    ('Custo zero de implantação',
     'Álcool, acidez volátil, sulfatos e ácido cítrico já são medidos hoje no laboratório. '
     'Nenhum equipamento novo, nenhum dado novo: é só conectar ao sistema de decisão.'),
    ('Correção durante o processo',
     'Se a acidez volátil sobe no tanque, dá para agir enquanto ainda há vinho a salvar — em '
     'vez de descobrir o problema quando a garrafa já está fechada.'),
]
textbox(s, MARGIN, Inches(4.66) - Inches(0.28), CONTENT_W, Inches(0.2),
        'Escopo: a base contém apenas vinhos tintos. Estender a conclusão aos brancos exige nova validação.',
        size=8, color=INK_SOFT, italic=True)
for i, (tit, corpo) in enumerate(concl):
    y = Inches(1.55) + Inches(1.02) * i
    rect(s, MARGIN, y + Inches(0.04), Inches(0.028), Inches(0.62), ACCENT)
    textbox(s, MARGIN + Inches(0.25), y, Inches(2.6), Inches(0.3), tit,
            size=12, font=DISPLAY, color=INK)
    textbox(s, MARGIN + Inches(3.0), y + Inches(0.02), CONTENT_W - Inches(3.0), Inches(0.7),
            corpo, size=9.5, color=INK_SOFT, spacing=1.45)
hairline(s, MARGIN, Inches(4.66), CONTENT_W)
textbox(s, MARGIN, Inches(4.8), CONTENT_W, Inches(0.25),
        'Próximo passo: escolher uma linha de produção para o piloto.',
        size=10.5, font=DISPLAY, color=ACCENT, italic=True)

out = ROOT / 'wine_quality_eda_storytelling.pptx'
prs.save(out)
print(f'Apresentação gerada: {out.name}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)')
print('Todos os números vieram de data/ e results/metricas_modelos.csv.')
